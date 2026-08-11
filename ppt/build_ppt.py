"""Builds slides.pptx from the content in PPT_Content.txt.
Run: python build_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

OUT_PATH = os.path.join(os.path.dirname(__file__), "slides.pptx")

TITLE_LAYOUT = 0
BULLET_LAYOUT = 1

ACCENT = RGBColor(0x1F, 0x4E, 0x79)


def add_bullet_slide(prs, title, bullets, font_size=18):
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_LAYOUT])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, (text, level) in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(font_size)
    return slide


def add_table_slide(prs, title, headers, rows, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(0.4), Inches(1.4), Inches(9.2), Inches(0.4 * n_rows)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    if note:
        note_box = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(0.6))
        tf = note_box.text_frame
        tf.text = note
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.italic = True
    return slide


def b(text, level=0):
    return (text, level)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1 - Title
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    slide.shapes.title.text = "Watermark-Guided Tamper Localization and Recovery"
    subtitle = slide.placeholders[1]
    subtitle.text = (
        "Self-embedding fragile watermarking that detects WHERE an image was altered—and rebuilds it\n"
        "Final Year Project — B.E./B.Tech CSE/IT\n"
        "Submitted by: (Student Name(s), Roll No.)\n"
        "Project Guide: (Guide Name, Designation)\n"
        "Department of CSE/IT, (College Name)\n"
        "Academic Year: (2026-27)"
    )

    # Slide 2 - Problem Statement
    add_bullet_slide(prs, "Problem Statement", [
        b("Editing an image is now a 10-second job — phone apps, one-click AI inpainting"),
        b("Tampered images drive real damage: misinformation, forged insurance/accident evidence, doctored ID and legal document scans, altered medical scans"),
        b("Courts and newsrooms need proof of integrity, not opinion"),
        b("Existing tools answer only one question: “Is this fake?”"),
        b("A single yes/no verdict is not actionable:"),
        b("Doesn't say which region was altered", 1),
        b("Doesn't say what was there before", 1),
        b("Destroyed evidence stays destroyed", 1),
        b("Gap addressed: detection + localization + RECOVERY of the original content"),
    ])

    # Slide 3 - Objectives
    add_bullet_slide(prs, "Objectives", [
        b("Embed a self-contained fragile watermark — the image carries its own backup, no external database"),
        b("Detect tampering at block level and output a precise tamper localization mask"),
        b("Recover an approximate version of tampered regions from redundant data stored elsewhere in the same image"),
        b("Keep the watermark imperceptible — target PSNR > 40 dB vs the original"),
        b("Build an automated evaluation pipeline: self-generated ground-truth masks across 4 tamper types, measured on precision/recall/F1/IoU and recovery PSNR/SSIM"),
        b("Achieve all of the above with classical, explainable signal processing — no training data, no GPU"),
    ])

    # Slide 4 - Literature Survey
    add_bullet_slide(prs, "Literature Survey", [
        b("Traditional passive forgery detection (copy-move, JPEG-artifact, noise/CFA inconsistency)"),
        b("No prior watermark needed, but heuristic, sensitive to post-processing, coarse localization, no recovery", 1),
        b("Deep-learning forgery / deepfake detectors (CNN, transformer, segmentation-based)"),
        b("Strong on benchmarks, but need large labelled datasets + GPU, generalize poorly, black-box scores", 1),
        b("Robust watermarking (DCT/DWT, spread-spectrum)"),
        b("Built to SURVIVE attacks for copyright proof — deliberately insensitive to tampering, unsuitable for authentication", 1),
        b("Fragile / semi-fragile watermarking"),
        b("Correct family for tamper detection; most schemes stop at detect-and-localize only; self-recovery variants underused in practice", 1),
        b("Consolidated gap: detection-only pipelines, weak explainability, training-data dependence, no restoration"),
    ])

    # Slide 5 - Proposed System Overview
    add_bullet_slide(prs, "Proposed System Overview", [
        b("Image treated as a grid of small 8x8 blocks"),
        b("Before distribution: every block compressed into a tiny “recovery code” hidden inside a distant partner block's LSBs; each block also carries a hash of itself"),
        b("Every block has a backup living somewhere far away in the same image"),
        b("After the image is shared/attacked: re-read every block, recompute its hash"),
        b("Match = authentic; Mismatch = tampered, marked on a binary mask", 1),
        b("For each tampered block: fetch its recovery code from its (surviving) partner and repaint the region"),
        b("Output: watermarked image, tamper mask, recovered image, plus full numeric report"),
        b("Plain language: the image carries its own insurance policy, and can heal itself"),
    ])

    # Slide 6 - Architecture (diagram placeholder + description)
    add_bullet_slide(prs, "System Architecture / Block Diagram", [
        b("(Diagram to insert: left-to-right pipeline)"),
        b("Stage 1: Original Image → Stage 2: Block Partitioning (8x8 grid)"),
        b("Stage 3: Recovery Data Generation (DCT low-freq coeffs / downsampled average → quantize)"),
        b("Stage 4: Key-Seeded Block Mapping (secret key → CSPRNG → single-cycle permutation m(i))"),
        b("Stage 5: LSB Embedding (partner's recovery bits + own index-bound hash) → Watermarked Image"),
        b("(Test path only) Tamper Simulation → Tampered Image + Ground-Truth Mask"),
        b("Stage 6: Detection (recompute hash, compare) → Tamper Localization Mask"),
        b("Stage 7: Recovery (read backup from partner, inverse-transform, repaint) → Recovered Image"),
        b("Stage 8: Evaluation (P/R/F1/IoU, recovery PSNR/SSIM, recoverability rate)"),
    ], font_size=16)

    # Slide 7 - Methodology Embedding
    add_bullet_slide(prs, "Methodology — Embedding", [
        b("Step 1: Load image, pad to multiple of 8, split into 8x8 blocks"),
        b("Step 2: Clear LSB plane — zero the 2 LSBs of every pixel (stable hash + deterministic embedding)"),
        b("Step 3: Recovery data per block — 2D DCT, keep low-frequency coefficients, quantize to a compact bitstring"),
        b("Step 4: Partner mapping — key-seeded single-cycle permutation m(i), minimum spatial separation enforced"),
        b("Reason: a local tamper must not destroy a block AND its own backup at once", 1),
        b("Step 5: Authentication hash — keyed hash bound to block index + image ID (defeats block-swap counterfeiting)"),
        b("Step 6: Embed — write partner's recovery bits + own hash tag into block's LSBs"),
        b("Output: watermarked image, visually identical, only LSBs changed"),
    ])

    # Slide 8 - Methodology Detection & Recovery
    add_bullet_slide(prs, "Methodology — Detection & Recovery", [
        b("Detection: recompute hash from current MSB content, compare to stored hash"),
        b("Match = authentic; Mismatch = tampered → binary tamper localization mask (with neighbourhood refinement)", 1),
        b("Recovery: for each tampered block, look up its partner"),
        b("Partner authentic → read recovery bits, inverse-transform, repaint block", 1),
        b("Partner ALSO tampered → marked explicitly UNRECOVERABLE, never silently faked", 1),
        b("Fully explainable: every flag traces to one named hash mismatch — no probability score, no black box"),
        b("Deterministic: same key + same image → identical result, every run"),
        b("Honest scope: built for a lossless pipeline (PNG) owned end-to-end, not for catching unknown external forgeries after re-compression"),
    ])

    # Slide 9 - Tools & Technologies (table)
    add_table_slide(
        prs, "Tools & Technologies Used",
        ["Tool", "Purpose"],
        [
            ["Python 3.x", "Core implementation language"],
            ["NumPy", "Block partitioning, LSB bit-plane manipulation, vectorized math"],
            ["OpenCV (cv2)", "Image I/O, DCT/IDCT, mask cleanup, inpainting fallback"],
            ["scikit-image", "PSNR and SSIM metric computation"],
            ["hashlib / hmac", "Keyed authentication tags bound to block index + image ID"],
            ["NumPy CSPRNG", "Secret-key-seeded block mapping"],
            ["Matplotlib", "Result plots and visual comparisons"],
            ["Streamlit", "Interactive demo UI — upload, embed, tamper, detect, recover live"],
            ["Pandas", "Aggregating metrics into result tables"],
        ],
        note="No TensorFlow/PyTorch, no dataset download, no GPU — runs fully offline, zero cost."
    )

    # Slide 10 - Implementation Snapshot
    add_bullet_slide(prs, "Implementation Snapshot", [
        b("(Figure to insert: 5-panel strip (a)-(e) side by side)"),
        b("(a) Original Image"),
        b("(b) Watermarked Image — visually identical, PSNR/SSIM caption"),
        b("(c) Tampered Image — object pasted in, or region removed/refilled"),
        b("(d) Detected Tamper Mask — white = flagged, overlay in red on tampered image"),
        b("(e) Recovered Image — repainted from partner-block backups, recovery PSNR caption"),
        b("(Screenshot to insert: Streamlit demo UI — key/block-size/tamper-type controls + live metrics)"),
    ])

    # Slide 11 - Results Imperceptibility (table) -- measured, 32-image corpus (8 USC-SIPI + 24 Kodak)
    add_table_slide(
        prs, "Results — Imperceptibility",
        ["Descriptor Variant", "PSNR (dB)", "SSIM"],
        [
            ["Variant A (DCT) — mean, 32 images", "43.17", "0.9824"],
            ["Variant B (mean-pooled) — mean, 32 images", "44.23", "0.9844"],
            ["Lena, Variant A / B", "43.22 / 44.31", "0.9814 / 0.9839"],
        ],
        note="Measured results: PSNR > 43 dB, SSIM > 0.97 on every image. Only 2 LSBs modified per pixel -> analytical reference point ~44.15 dB."
    )

    # Slide 12 - Results Detection & Localization (table) -- measured, block-level, 864 tamper trials
    add_table_slide(
        prs, "Results — Detection & Localization",
        ["Tamper Type", "Block Prec.", "Block Recall", "Block F1", "Block IoU"],
        [
            ["Copy-paste splicing", "1.0000", "1.0000", "1.0000", "1.0000"],
            ["Object removal / inpainting", "1.0000", "1.0000", "1.0000", "1.0000"],
            ["Region crop & refill", "1.0000", "1.0000", "1.0000", "1.0000"],
            ["Noise-block corruption", "1.0000", "1.0000", "1.0000", "1.0000"],
            ["Overall (864 trials)", "1.0000", "1.0000", "1.0000", "1.0000"],
        ],
        note="Ground truth is exact (self-generated). Block metrics saturated near 1.0; pixel-level precision mean 0.9489 (gap is block-grid quantization, not false alarms). Security fix: binding the recovery descriptor into the tag eliminated the MSB-preserved miss category (was 0.143 blocks/row, mostly object removal) -- recall is now 1.0000 on every class. Null condition: 0 false positives / 1,802,240 block checks (rule-of-three 95% bound 1.66e-6)."
    )

    # Slide 13 - Results Recovery Quality (table) -- measured, Variant A, by tamper ratio
    # Restructured from tamper-TYPE to tamper-RATIO: that's the axis the most interesting
    # measured finding (flat in-region PSNR vs. collapsing whole-image PSNR) actually varies over.
    add_table_slide(
        prs, "Results — Recovery Quality",
        ["Tamper Ratio", "Recoverability (rho)", "In-Region PSNR (dB)", "Whole-Image PSNR (dB)"],
        [
            ["10%", "0.9590", "28.96", "33.79"],
            ["25%", "0.8036", "28.44", "23.87"],
            ["50%", "0.5335", "28.38", "17.39"],
        ],
        note="Key finding: in-region PSNR stays flat as tamper ratio grows; whole-image PSNR collapses. Coverage (rho), not descriptor fidelity, drives whole-image quality. Security/perf trade-off: removing the block-mapping structural leak cost ~1.5 points of rho at every ratio (leaky map 0.9736/0.8353/0.5469 vs flat map 0.9590/0.8036/0.5335) -- the price of removing a bias an attacker with no key could exploit. At 50%, well below Korus & Dziech's 37 dB (reference-sharing degrades gracefully; our 1:1 mapping does not -- see Future Scope)."
    )

    # Slide 14 - Comparison table
    # Includes the learned proactive-watermarking family (EditGuard/DeepMark/RecoverMark) so the
    # table does NOT imply we uniquely recover -- they also localize and recover. Our defensible
    # axis vs. that family is determinism, zero training, zero GPU, and full auditability.
    add_table_slide(
        prs, "Comparison with Existing Approaches",
        ["Capability", "Passive Detection", "Deep-Learning Detectors", "Fragile WM (detect-only)",
         "Learned Proactive WM (EditGuard/DeepMark/RecoverMark)", "Proposed System"],
        [
            ["Detects tampering", "Partial", "Yes", "Yes", "Yes", "Yes"],
            ["Localizes tampering", "Coarse", "Yes", "Yes", "Yes (learned)", "Yes (block-precise)"],
            ["Recovers content", "No", "No", "No", "Yes (learned)", "Yes (deterministic)"],
            ["Needs training data", "No", "Yes, large+labelled", "No", "Yes, large", "No"],
            ["Needs GPU", "No", "Usually", "No", "Yes (train+infer)", "No"],
            ["Explainable decision", "Partly", "No (black-box)", "Yes", "No (network output)", "Yes (traces to 1 hash)"],
            ["Deterministic", "Mostly", "No", "Yes", "No", "Yes"],
            ["Requires prior watermark", "No", "No", "Yes", "Yes", "Yes (accepted trade-off)"],
        ],
        note="Honest framing: recent learned proactive watermarks ALSO localize and recover -- we do not claim unique recovery. Our defensible axis vs. that family: zero training, zero GPU, deterministic, and fully auditable (every decision traces to one recomputed keyed hash)."
    )

    # Slide 15 - Applications
    add_bullet_slide(prs, "Applications", [
        b("Digital forensic evidence — CCTV, crime-scene/accident photos watermarked at capture"),
        b("Journalism & media authentication — news agencies watermark on ingest"),
        b("Medical imaging integrity — X-ray/MRI/CT in PACS, altered diagnostic regions flagged and approximately restored"),
        b("Legal & official document scans — cheques, contracts, mark sheets, ID cards, land records"),
        b("Insurance claim verification — damage photos proven unedited before settlement"),
        b("Government / e-governance archives — long-term integrity of scanned public records"),
        b("Secure camera / device firmware — self-authenticating images at the sensor level"),
    ])

    # Slide 16 - Limitations & Future Scope
    add_bullet_slide(prs, "Limitations & Future Scope", [
        b("Limitations:"),
        b("Fragile to benign processing — JPEG re-compression, resizing destroy the watermark; assumes lossless storage", 1),
        b("Block-size trade-off — smaller blocks = sharper localization but less recovery capacity", 1),
        b("Recovery is approximate, not pixel-perfect", 1),
        b("Large-area tampering degrades recovery — block + partner both destroyed = unrecoverable", 1),
        b("Security depends fully on secret-key confidentiality", 1),
        b("Future Scope:"),
        b("Semi-fragile/robust embedding to tolerate mild JPEG compression", 1),
        b("Reference-sharing / redundant backups to survive large-area attacks", 1),
        b("Small unsupervised VQ codebook for smarter recovery encoding (no GPU, no labels)", 1),
        b("Extension to video; camera-side SDK / mobile app integration", 1),
    ])

    # Slide 17 - Conclusion
    add_bullet_slide(prs, "Conclusion", [
        b("Designed and implemented a block-based self-embedding fragile watermarking system"),
        b("Detects tampering, localizes it to individual 8x8 blocks, and recovers the altered content from key-mapped partner-block backups"),
        b("Achieved imperceptible embedding (PSNR 43.17-44.23 dB across descriptor variants), block-level localization precision 1.0000 / recall 1.0000, and 28.4-29.0 dB in-region recovery PSNR that stays flat across tamper ratio even as whole-image PSNR falls from 33.79 to 17.39 dB"),
        b("Evaluated across four tamper types with exact self-generated ground truth"),
        b("Delivered with classical, deterministic Python/OpenCV signal processing — explainable, reproducible, offline, zero-cost"),
        b("Goes one decisive step beyond conventional forgery detectors: they only flag tampering, this restores what was lost"),
    ])

    # Slide 18 - References
    add_bullet_slide(prs, "References", [
        b("(Reference list to insert here — must match the IEEE paper's reference list)"),
        b("Coverage: fragile & self-recovery watermarking, DCT-based image authentication, LSB steganography, image forgery detection surveys, deep-learning forgery detection, PSNR/SSIM metrics, OpenCV/scikit-image docs"),
    ])

    # Slide 19 - Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    slide.shapes.title.text = "Thank You"
    subtitle = slide.placeholders[1]
    subtitle.text = (
        "Questions & Discussion\n"
        "Live demo available on request — upload an image, tamper it, watch it heal\n"
        "Contact: (email) | Repository: (link)\n"
        "Guide: (Guide Name) | Department of CSE/IT, (College Name)"
    )

    prs.save(OUT_PATH)
    print(f"Saved {OUT_PATH} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
